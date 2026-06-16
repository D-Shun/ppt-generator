from __future__ import annotations

import argparse
import os
import shutil
import sys
import traceback
from copy import deepcopy
from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple

from openpyxl import load_workbook
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from PyQt5.QtCore import QThread, Qt, pyqtSignal
from PyQt5.QtGui import QFont
from PyQt5.QtWidgets import (
    QApplication,
    QFileDialog,
    QFormLayout,
    QFrame,
    QGridLayout,
    QHBoxLayout,
    QGroupBox,
    QLabel,
    QLineEdit,
    QMainWindow,
    QMessageBox,
    QPushButton,
    QComboBox,
    QPlainTextEdit,
    QProgressBar,
    QVBoxLayout,
    QWidget,
)


ROOT = Path.cwd()
DEFAULT_EXCEL = next((p for p in ROOT.glob("*.xlsx") if not p.name.startswith("~$")), None)
DEFAULT_PPT = next(ROOT.glob("*.pptx"), None)
DEFAULT_OUTPUT = ROOT / "output" / "测试生成结果.pptx"
DEFAULT_OUTPUT_DIR = ROOT / "output"
TEMP_DIR = ROOT / "temp_images"


@dataclass
class RowRecord:
    excel_row: int
    values: Dict[str, str]
    image_path: Optional[Path] = None

    @property
    def has_image(self) -> bool:
        return self.image_path is not None and self.image_path.exists()


def normalize_value(value) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def display_path(path: Path) -> str:
    try:
        return str(path.resolve())
    except Exception:
        return str(path)


def safe_filename(text: str) -> str:
    cleaned = "".join(ch for ch in text if ch not in '<>:"/\\|?*')
    cleaned = cleaned.strip().strip(".")
    return cleaned or "输出文件"


def resolve_output_target(output_text: str) -> Path:
    path = Path(output_text.strip())
    if path.suffix.lower() == ".pptx":
        return path if path.is_absolute() else ROOT / path
    folder = path if path.is_absolute() else ROOT / path
    return folder / f"{safe_filename(DEFAULT_OUTPUT.stem)}.pptx"


def resolve_path_text(text: str) -> Path:
    path = Path(text.strip())
    if not path.is_absolute():
        path = ROOT / path
    return path


def detect_image_ext(data: bytes, fallback: str = ".png") -> str:
    try:
        with PILImage.open(BytesIO(data)) as im:
            fmt = (im.format or "").lower()
            if fmt == "jpeg":
                return ".jpg"
            if fmt:
                return f".{fmt}"
    except Exception:
        pass
    return fallback


def clean_temp_dir(path: Path) -> None:
    if path.exists():
        for item in path.iterdir():
            try:
                if item.is_file() or item.is_symlink():
                    item.unlink()
                elif item.is_dir():
                    shutil.rmtree(item)
            except Exception:
                pass
    else:
        path.mkdir(parents=True, exist_ok=True)


def extract_excel_rows_and_images(
    excel_path: Path,
    sheet_name: str,
    temp_dir: Path,
) -> Tuple[List[str], List[RowRecord], Dict[str, int]]:
    wb = load_workbook(excel_path, data_only=True)
    if sheet_name not in wb.sheetnames:
        raise ValueError(f"工作表不存在: {sheet_name}")

    ws = wb[sheet_name]
    headers = [normalize_value(c.value).strip() for c in ws[1]]
    headers = [h for h in headers if h]
    if not headers:
        raise ValueError("Excel 第一行没有找到表头")

    clean_temp_dir(temp_dir)

    image_by_row: Dict[int, Path] = {}
    embedded_images = getattr(ws, "_images", []) or []
    image_count = 0
    for idx, img in enumerate(embedded_images, start=1):
        try:
            anchor = img.anchor
            excel_row = int(anchor._from.row) + 1
            data = img._data()
            ext = detect_image_ext(data, ".png")
            image_path = temp_dir / f"row_{excel_row}_{idx}{ext}"
            image_path.write_bytes(data)
            image_by_row.setdefault(excel_row, image_path)
            image_count += 1
        except Exception:
            continue

    records: List[RowRecord] = []
    max_row = ws.max_row or 1
    for row_idx in range(2, max_row + 1):
        values: Dict[str, str] = {}
        has_any = False
        for col_idx, header in enumerate(headers, start=1):
            cell_value = normalize_value(ws.cell(row=row_idx, column=col_idx).value)
            values[header] = cell_value
            if cell_value:
                has_any = True
        if not has_any and row_idx not in image_by_row:
            continue
        records.append(RowRecord(excel_row=row_idx, values=values, image_path=image_by_row.get(row_idx)))

    stats = {
        "rows": len(records),
        "images": image_count,
        "rows_with_image": sum(1 for r in records if r.has_image),
    }
    return headers, records, stats


def set_textbox_style(
    shape,
    *,
    font_name: str = "Microsoft YaHei",
    font_size: int = 16,
    bold: bool = False,
    color: Tuple[int, int, int] = (34, 34, 34),
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
) -> None:
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = valign
    for p in text_frame.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)


def add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: int,
    bold: bool = False,
    color: Tuple[int, int, int] = (34, 34, 34),
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
    fill: Optional[Tuple[int, int, int]] = None,
    line: Optional[Tuple[int, int, int]] = None,
) -> object:
    shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = RGBColor(*line)
    else:
        shape.line.fill.background()
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    set_textbox_style(shape, font_size=font_size, bold=bold, color=color, align=align, valign=valign)
    return shape


def add_fitted_picture(slide, image_path: Path, left: int, top: int, width: int, height: int) -> None:
    with PILImage.open(image_path) as im:
        img_w, img_h = im.size
    if img_w <= 0 or img_h <= 0:
        return
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio >= box_ratio:
        pic_w = width
        pic_h = max(1, int(width / img_ratio))
    else:
        pic_h = height
        pic_w = max(1, int(height * img_ratio))
    pic_left = left + int((width - pic_w) / 2)
    pic_top = top + int((height - pic_h) / 2)
    slide.shapes.add_picture(str(image_path), Emu(pic_left), Emu(pic_top), Emu(pic_w), Emu(pic_h))


def unique_output_path(path: Path) -> Path:
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    return path.with_name(f"{path.stem}-{stamp}{path.suffix or '.pptx'}")


def save_presentation_with_fallback(prs: Presentation, output_path: Path, log: Callable[[str], None]) -> Path:
    try:
        prs.save(str(output_path))
        return output_path
    except PermissionError:
        fallback = unique_output_path(output_path)
        log(f"目标文件可能正在被打开，已自动改名输出: {fallback.name}")
        prs.save(str(fallback))
        return fallback


def add_placeholder(slide, left: int, top: int, width: int, height: int, text: str = "无图片") -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(243, 243, 243)
    shape.line.color.rgb = RGBColor(210, 210, 210)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(140, 140, 140)


def clone_slide(prs: Presentation, template_slide):
    blank_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    for shape in list(new_slide.shapes):
        shape.element.getparent().remove(shape.element)

    for shape in template_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

    for rel in template_slide.part.rels.values():
        if rel.is_external:
            new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_slide.part.relate_to(rel.target_part, rel.reltype)

    return new_slide


def remove_slide(prs: Presentation, slide) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = None
    for candidate in slide_id_list:
        if prs.part.related_part(candidate.rId) == slide.part:
            slide_id = candidate
            break
    if slide_id is not None:
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def replace_shape_text(shape, text: str) -> None:
    if not hasattr(shape, "text_frame"):
        return
    text_frame = shape.text_frame
    lines = str(text).splitlines() or [""]
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = str(text)
        return

    for idx, line in enumerate(lines):
        if idx >= len(paragraphs):
            paragraph = text_frame.add_paragraph()
            # New paragraphs are rare for this template. Keep them modest if the
            # Excel text has more lines than the template sample.
            if paragraphs and paragraphs[-1].runs:
                source_font = paragraphs[-1].runs[0].font
                run = paragraph.add_run()
                run.font.name = source_font.name
                run.font.size = source_font.size
                run.font.bold = source_font.bold
            else:
                run = paragraph.add_run()
            paragraphs.append(paragraph)
        paragraph = paragraphs[idx]
        if paragraph.runs:
            paragraph.runs[0].text = line
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            run = paragraph.add_run()
            run.text = line

    for paragraph in paragraphs[len(lines):]:
        for run in paragraph.runs:
            run.text = ""


def shape_plain_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return (shape.text or "").replace("\n", "").strip()


def classify_template_shape(shape) -> Optional[str]:
    if not hasattr(shape, "text_frame"):
        return None
    text = shape_plain_text(shape)
    name = getattr(shape, "name", "")
    top = int(getattr(shape, "top", 0))

    if "D098821" in text:
        return "商品编码"
    if "商品名称" in text or (name == "标题 1" and top > 1100000):
        return "商品名称"
    if "如水" in text or name == "文本框 3":
        return "商品规格"
    if text in {"顺丰/申通", "顺丰", "申通"} or name == "文本框 5":
        return "快递公司"
    if text == "2" or name == "文本框 4":
        return "包裹数量"
    if text == "北京" or name == "文本框 6":
        return "发货地"
    return None


def is_picture_shape(shape) -> bool:
    try:
        return int(shape.shape_type) == 13
    except Exception:
        return False


def find_template_picture(slide):
    pictures = [shape for shape in slide.shapes if is_picture_shape(shape)]
    if not pictures:
        return None
    return max(pictures, key=lambda sh: int(sh.width) * int(sh.height))


def remove_shape(shape) -> None:
    shape.element.getparent().remove(shape.element)


def replace_copied_slide_content(slide, record: RowRecord) -> None:
    for shape in list(slide.shapes):
        field = classify_template_shape(shape)
        if field:
            replace_shape_text(shape, record.values.get(field, ""))

    picture_shape = find_template_picture(slide)
    if picture_shape is None:
        return

    left = int(picture_shape.left)
    top = int(picture_shape.top)
    width = int(picture_shape.width)
    height = int(picture_shape.height)
    remove_shape(picture_shape)

    if record.image_path and record.image_path.exists():
        add_fitted_picture(slide, record.image_path, left, top, width, height)


def build_slide(slide, record: RowRecord, slide_size: Tuple[int, int]) -> None:
    slide_width, slide_height = slide_size

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(slide_width), Emu(slide_height))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    bg.element.getparent().remove(bg.element)
    slide.shapes._spTree.insert(2, bg.element)

    # Outer frame and header accents
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(420000),
        Emu(1100000),
        Emu(slide_width - 840000),
        Emu(slide_height - 1700000),
    )
    frame.fill.background()
    frame.line.color.rgb = RGBColor(237, 103, 38)
    frame.line.width = Pt(3)

    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(420000),
        Emu(420000),
        Emu(slide_width - 840000),
        Emu(300000),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(250, 111, 38)
    top_bar.line.fill.background()

    inner_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(420000),
        Emu(1060000),
        Emu(slide_width - 840000),
        Emu(70000),
    )
    inner_bar.fill.solid()
    inner_bar.fill.fore_color.rgb = RGBColor(255, 202, 58)
    inner_bar.line.fill.background()

    code = record.values.get("商品编码", "")
    name = record.values.get("商品名称", "")
    spec = record.values.get("商品规格", "")
    ship = record.values.get("快递公司", "")
    qty = record.values.get("包裹数量", "")
    origin = record.values.get("发货地", "")

    add_textbox(
        slide,
        930000,
        920000,
        3100000,
        380000,
        code,
        font_size=14,
        bold=True,
        color=(45, 45, 45),
    )
    add_textbox(
        slide,
        970000,
        1210000,
        8650000,
        600000,
        name,
        font_size=24,
        bold=True,
        color=(26, 26, 26),
    )
    add_textbox(
        slide,
        1080000,
        1900000,
        4300000,
        3850000,
        spec,
        font_size=15,
        bold=False,
        color=(34, 34, 34),
        valign=MSO_VERTICAL_ANCHOR.TOP,
    )

    image_left = 5850000
    image_top = 1780000
    image_width = 3600000
    image_height = 3600000
    if record.image_path and record.image_path.exists():
        add_fitted_picture(slide, record.image_path, image_left, image_top, image_width, image_height)
    else:
        add_placeholder(slide, image_left, image_top, image_width, image_height)

    footer_y = 6720000
    add_textbox(
        slide,
        6040000,
        footer_y,
        940000,
        320000,
        ship,
        font_size=12,
        bold=False,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
        color=(70, 70, 70),
    )
    add_textbox(
        slide,
        7600000,
        footer_y,
        940000,
        320000,
        qty,
        font_size=12,
        bold=False,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
        color=(70, 70, 70),
    )
    add_textbox(
        slide,
        9180000,
        footer_y,
        940000,
        320000,
        origin,
        font_size=12,
        bold=False,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
        color=(70, 70, 70),
    )

    # Bottom hint line to make the test output easier to identify.
    add_textbox(
        slide,
        420000,
        slide_height - 450000,
        2600000,
        220000,
        f"第 {record.excel_row} 行",
        font_size=10,
        color=(200, 83, 34),
    )


def generate_ppt(
    excel_path: Path,
    ppt_path: Path,
    output_path: Path,
    sheet_name: Optional[str],
    log: Callable[[str], None] = lambda *_: None,
    progress: Optional[Callable[[int, int], None]] = None,
) -> Tuple[Dict[str, int], Path]:
    if not excel_path.exists():
        raise FileNotFoundError(f"Excel 文件不存在: {excel_path}")
    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT 模板不存在: {ppt_path}")

    prs = Presentation(str(ppt_path))
    if not prs.slides:
        raise ValueError("PPT 模板里没有幻灯片")
    template_slide = prs.slides[0]

    wb = load_workbook(excel_path, data_only=True)
    target_sheet = sheet_name or wb.sheetnames[0]
    headers, records, stats = extract_excel_rows_and_images(excel_path, target_sheet, TEMP_DIR)
    if not records:
        raise ValueError("Excel 里没有可生成的数据行")

    log(f"读取工作表: {target_sheet}")
    log(f"表头: {' | '.join(headers)}")
    log(f"数据行: {len(records)}，嵌入图片: {stats['images']}，带图行: {stats['rows_with_image']}")
    log("生成方式: 复制原 PPT 模板页，每一行 Excel 数据生成一页")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    total = len(records)
    for idx, record in enumerate(records, start=1):
        slide = clone_slide(prs, template_slide)
        replace_copied_slide_content(slide, record)
        if progress:
            progress(idx, total)
        has_img = "有图" if record.has_image else "无图"
        log(f"已生成第 {idx}/{total} 页: Excel 第 {record.excel_row} 行 ({has_img})")

    remove_slide(prs, template_slide)
    saved_path = save_presentation_with_fallback(prs, output_path, log)
    log(f"输出完成: {saved_path}")
    return stats, saved_path


class GenerateWorker(QThread):
    log_signal = pyqtSignal(str)
    progress_signal = pyqtSignal(int, int)
    done_signal = pyqtSignal(str)
    error_signal = pyqtSignal(str)

    def __init__(self, excel_path: Path, ppt_path: Path, output_path: Path, sheet_name: str):
        super().__init__()
        self.excel_path = excel_path
        self.ppt_path = ppt_path
        self.output_path = output_path
        self.sheet_name = sheet_name

    def run(self) -> None:
        try:
            _, saved_path = generate_ppt(
                self.excel_path,
                self.ppt_path,
                self.output_path,
                self.sheet_name,
                log=self.log_signal.emit,
                progress=lambda cur, total: self.progress_signal.emit(cur, total),
            )
            self.done_signal.emit(str(saved_path))
        except Exception:
            self.error_signal.emit(traceback.format_exc())


class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Excel 批量生成 PPT 测试版")
        self.resize(1280, 820)
        self.setMinimumSize(1160, 720)

        self.headers: List[str] = []
        self.records: List[RowRecord] = []
        self.stats: Dict[str, int] = {}
        self.worker: Optional[GenerateWorker] = None

        self.excel_edit = QLineEdit()
        self.ppt_edit = QLineEdit()
        self.output_edit = QLineEdit()
        self.sheet_combo = QComboBox()
        self.summary_label = QLabel("未加载")
        self.progress_bar = QProgressBar()
        self.log_box = QPlainTextEdit()

        self._build_ui()
        self._load_defaults()

    def _build_ui(self) -> None:
        font = QFont("Microsoft YaHei", 10)
        QApplication.instance().setFont(font)

        central = QWidget()
        self.setCentralWidget(central)

        root_layout = QGridLayout(central)
        root_layout.setContentsMargins(16, 16, 16, 16)
        root_layout.setHorizontalSpacing(14)
        root_layout.setVerticalSpacing(12)

        header = QWidget()
        header_layout = QVBoxLayout(header)
        header_layout.setContentsMargins(0, 0, 0, 0)
        header_layout.setSpacing(4)
        title = QLabel("Excel 批量生成 PPT")
        title.setObjectName("pageTitle")
        subtitle = QLabel("选择 Excel、模板和输出位置，批量生成商品 PPT。")
        subtitle.setObjectName("pageSubtitle")
        header_layout.addWidget(title)
        header_layout.addWidget(subtitle)
        root_layout.addWidget(header, 0, 0, 1, 2)

        left_panel = QWidget()
        left_panel.setMinimumWidth(340)
        left_panel.setMaximumWidth(390)
        left_layout = QVBoxLayout(left_panel)
        left_layout.setContentsMargins(0, 0, 0, 0)
        left_layout.setSpacing(12)

        input_box = QGroupBox("输入")
        input_box.setObjectName("panelBox")
        input_form = QFormLayout(input_box)
        input_form.setLabelAlignment(Qt.AlignLeft)
        input_form.setFormAlignment(Qt.AlignTop)
        input_form.setHorizontalSpacing(10)
        input_form.setVerticalSpacing(10)

        input_form.addRow(self._path_row(self.excel_edit, "Excel", self.browse_excel))
        input_form.addRow(self._path_row(self.ppt_edit, "PPT", self.browse_ppt))
        input_form.addRow(self._path_row(self.output_edit, "输出", self.browse_output))
        input_form.addRow("工作表", self.sheet_combo)

        self.sheet_combo.currentTextChanged.connect(self.refresh_preview)

        left_layout.addWidget(input_box)

        stat_box = QGroupBox("摘要")
        stat_box.setObjectName("panelBox")
        stat_layout = QVBoxLayout(stat_box)
        self.summary_label.setWordWrap(True)
        self.summary_label.setStyleSheet("color:#39424e;")
        self.summary_label.setMinimumHeight(96)
        self.summary_label.setAlignment(Qt.AlignTop | Qt.AlignLeft)
        stat_layout.addWidget(self.summary_label)
        left_layout.addWidget(stat_box)

        action_box = QGroupBox("操作")
        action_box.setObjectName("panelBox")
        action_layout = QVBoxLayout(action_box)
        action_layout.setSpacing(10)
        self.generate_btn = QPushButton("生成测试 PPT")
        self.generate_btn.clicked.connect(self.start_generation)
        self.generate_btn.setMinimumHeight(42)
        action_layout.addWidget(self.generate_btn)

        self.open_output_btn = QPushButton("打开输出目录")
        self.open_output_btn.clicked.connect(self.open_output_folder)
        self.open_output_btn.setMinimumHeight(38)
        action_layout.addWidget(self.open_output_btn)

        self.reload_btn = QPushButton("重新读取 Excel")
        self.reload_btn.clicked.connect(self.refresh_preview)
        self.reload_btn.setMinimumHeight(38)
        action_layout.addWidget(self.reload_btn)

        left_layout.addWidget(action_box)
        left_layout.addStretch(1)

        log_box = QGroupBox("日志")
        log_box.setObjectName("panelBox")
        log_layout = QVBoxLayout(log_box)
        self.log_box.setReadOnly(True)
        self.log_box.setFrameShape(QFrame.NoFrame)
        self.log_box.setPlaceholderText("运行日志会显示在这里")
        self.log_box.setMinimumHeight(620)
        log_layout.addWidget(self.log_box)

        root_layout.addWidget(left_panel, 1, 0)
        root_layout.addWidget(log_box, 1, 1)
        root_layout.setColumnStretch(0, 0)
        root_layout.setColumnStretch(1, 1)
        root_layout.setRowStretch(1, 1)

        self.progress_bar.setMinimum(0)
        self.progress_bar.setMaximum(100)
        self.progress_bar.setValue(0)
        self.statusBar().addPermanentWidget(self.progress_bar, 1)

        self.setStyleSheet(
            """
            QMainWindow { background: #f4f6f9; }
            QLabel#pageTitle {
                font-size: 22px;
                font-weight: 700;
                color: #111827;
            }
            QLabel#pageSubtitle {
                font-size: 12px;
                color: #667085;
            }
            QGroupBox {
                border: 1px solid #d9dee7;
                border-radius: 6px;
                margin-top: 12px;
                background: white;
                font-weight: 600;
            }
            QGroupBox::title {
                subcontrol-origin: margin;
                left: 12px;
                padding: 0 6px 0 6px;
                color: #1f2937;
            }
            QGroupBox#panelBox {
                margin-top: 14px;
            }
            QLineEdit, QComboBox, QPlainTextEdit {
                border: 1px solid #d0d7e2;
                border-radius: 4px;
                background: white;
                padding: 6px 8px;
            }
            QPushButton {
                border: 1px solid #cfd6e2;
                border-radius: 4px;
                padding: 8px 10px;
                background: #ffffff;
            }
            QPushButton:hover { background: #f5f7fa; }
            QPushButton:pressed { background: #e8edf4; }
            """
        )

    def _path_row(self, edit: QLineEdit, label: str, handler: Callable[[], None]) -> QWidget:
        container = QWidget()
        layout = QGridLayout(container)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setHorizontalSpacing(10)
        layout.setColumnStretch(0, 1)
        edit.setPlaceholderText(f"请选择{label}")
        edit.setMinimumHeight(34)
        edit.setToolTip("当前路径")
        layout.addWidget(edit, 0, 0)
        btn = QPushButton(f"选择{label}")
        btn.setMinimumHeight(34)
        btn.clicked.connect(handler)
        layout.addWidget(btn, 0, 1)
        return container

    def _load_defaults(self) -> None:
        if DEFAULT_EXCEL:
            self.excel_edit.setText(display_path(DEFAULT_EXCEL))
            self.excel_edit.setToolTip(str(DEFAULT_EXCEL.resolve()))
        if DEFAULT_PPT:
            self.ppt_edit.setText(display_path(DEFAULT_PPT))
            self.ppt_edit.setToolTip(str(DEFAULT_PPT.resolve()))
        self.output_edit.setText(display_path(DEFAULT_OUTPUT_DIR))
        self.output_edit.setToolTip(str(DEFAULT_OUTPUT_DIR.resolve()))
        if self.excel_edit.text():
            self.refresh_preview()

    def append_log(self, text: str) -> None:
        self.log_box.appendPlainText(text)

    def browse_excel(self) -> None:
        path, _ = QFileDialog.getOpenFileName(self, "选择 Excel 文件", str(ROOT), "Excel Files (*.xlsx *.xlsm)")
        if path:
            p = Path(path)
            self.excel_edit.setText(display_path(p))
            self.excel_edit.setToolTip(str(p.resolve()))
            self.refresh_preview()

    def browse_ppt(self) -> None:
        path, _ = QFileDialog.getOpenFileName(self, "选择 PPT 模板", str(ROOT), "PowerPoint Files (*.pptx)")
        if path:
            p = Path(path)
            self.ppt_edit.setText(display_path(p))
            self.ppt_edit.setToolTip(str(p.resolve()))

    def browse_output(self) -> None:
        path, _ = QFileDialog.getSaveFileName(self, "选择输出文件", self.output_edit.text() or str(DEFAULT_OUTPUT), "PowerPoint Files (*.pptx)")
        if path:
            if not path.lower().endswith(".pptx"):
                path += ".pptx"
            p = Path(path)
            self.output_edit.setText(display_path(p))
            self.output_edit.setToolTip(str(p.resolve()))

    def refresh_preview(self) -> None:
        excel_path = resolve_path_text(self.excel_edit.text())
        if not excel_path.exists():
            self.summary_label.setText("请选择有效的 Excel 文件")
            self.sheet_combo.blockSignals(True)
            self.sheet_combo.clear()
            self.sheet_combo.blockSignals(False)
            return

        try:
            wb = load_workbook(excel_path, data_only=True)
            sheets = wb.sheetnames
            self.sheet_combo.blockSignals(True)
            current = self.sheet_combo.currentText()
            self.sheet_combo.clear()
            self.sheet_combo.addItems(sheets)
            if current in sheets:
                self.sheet_combo.setCurrentText(current)
            self.sheet_combo.blockSignals(False)
            target_sheet = self.sheet_combo.currentText() or sheets[0]
            headers, records, stats = extract_excel_rows_and_images(excel_path, target_sheet, TEMP_DIR)
            self.headers = headers
            self.records = records
            self.stats = stats
            self._update_summary()
            self.append_log(f"已读取: {excel_path.name} / {target_sheet}")
        except Exception as exc:
            self.summary_label.setText(f"读取失败: {exc}")
            self.append_log(traceback.format_exc())

    def _update_summary(self) -> None:
        if not self.records:
            self.summary_label.setText("没有可生成的数据")
            return
        self.summary_label.setText(
            f"字段数: {len(self.headers)}\n"
            f"数据行: {self.stats.get('rows', 0)}\n"
            f"嵌入图片: {self.stats.get('images', 0)}\n"
            f"带图行: {self.stats.get('rows_with_image', 0)}"
        )

    def start_generation(self) -> None:
        excel_path = resolve_path_text(self.excel_edit.text())
        ppt_path = resolve_path_text(self.ppt_edit.text())
        output_path = resolve_output_target(self.output_edit.text())
        sheet_name = self.sheet_combo.currentText().strip()

        if not excel_path.exists():
            QMessageBox.warning(self, "提示", "请先选择有效的 Excel 文件")
            return
        if not ppt_path.exists():
            QMessageBox.warning(self, "提示", "请先选择有效的 PPT 模板")
            return
        if not sheet_name:
            QMessageBox.warning(self, "提示", "请先选择工作表")
            return

        self.generate_btn.setEnabled(False)
        self.progress_bar.setValue(0)
        self.append_log("开始生成...")
        self.worker = GenerateWorker(excel_path, ppt_path, output_path, sheet_name)
        self.worker.log_signal.connect(self.append_log)
        self.worker.progress_signal.connect(self.on_progress)
        self.worker.done_signal.connect(self.on_done)
        self.worker.error_signal.connect(self.on_error)
        self.worker.finished.connect(lambda: self.generate_btn.setEnabled(True))
        self.worker.start()

    def on_progress(self, cur: int, total: int) -> None:
        value = int(cur * 100 / total) if total else 0
        self.progress_bar.setValue(value)
        self.statusBar().showMessage(f"生成中 {cur}/{total}")

    def on_done(self, output_file: str) -> None:
        self.progress_bar.setValue(100)
        self.statusBar().showMessage("生成完成", 5000)
        self.append_log(f"完成: {output_file}")
        QMessageBox.information(self, "完成", f"PPT 已生成:\n{output_file}")

    def on_error(self, error_text: str) -> None:
        self.progress_bar.setValue(0)
        self.append_log(error_text)
        QMessageBox.critical(self, "生成失败", error_text)

    def open_output_folder(self) -> None:
        output_path = resolve_path_text(self.output_edit.text() or str(DEFAULT_OUTPUT))
        folder = output_path.parent
        folder.mkdir(parents=True, exist_ok=True)
        try:
            os.startfile(str(folder))
        except Exception as exc:
            QMessageBox.warning(self, "提示", f"无法打开输出目录: {exc}")


class CleanMainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Excel 批量生成 PPT 测试版")
        self.resize(1220, 780)
        self.setMinimumSize(1080, 680)

        self.headers: List[str] = []
        self.records: List[RowRecord] = []
        self.stats: Dict[str, int] = {}
        self.worker: Optional[GenerateWorker] = None

        self.excel_edit = QLineEdit()
        self.ppt_edit = QLineEdit()
        self.output_edit = QLineEdit()
        self.sheet_combo = QComboBox()
        self.log_box = QPlainTextEdit()
        self.progress_bar = QProgressBar()

        self.fields_value = QLabel("-")
        self.rows_value = QLabel("-")
        self.images_value = QLabel("-")
        self.ready_label = QLabel("等待读取 Excel")

        self._build_ui()
        self._load_defaults()

    def _build_ui(self) -> None:
        QApplication.instance().setFont(QFont("Microsoft YaHei", 10))

        central = QWidget()
        central.setObjectName("root")
        self.setCentralWidget(central)

        root = QVBoxLayout(central)
        root.setContentsMargins(22, 20, 22, 18)
        root.setSpacing(16)

        header = QWidget()
        header_layout = QHBoxLayout(header)
        header_layout.setContentsMargins(0, 0, 0, 0)
        header_layout.setSpacing(16)

        title_box = QWidget()
        title_layout = QVBoxLayout(title_box)
        title_layout.setContentsMargins(0, 0, 0, 0)
        title_layout.setSpacing(4)

        title = QLabel("Excel 批量生成 PPT")
        title.setObjectName("title")
        subtitle = QLabel("选择表格、模板和输出文件，一键生成商品介绍 PPT。")
        subtitle.setObjectName("subtitle")
        title_layout.addWidget(title)
        title_layout.addWidget(subtitle)

        self.ready_label.setObjectName("statusPill")
        self.ready_label.setAlignment(Qt.AlignCenter)
        self.ready_label.setMinimumWidth(150)

        header_layout.addWidget(title_box, 1)
        header_layout.addWidget(self.ready_label, 0, Qt.AlignTop)
        root.addWidget(header)

        body = QHBoxLayout()
        body.setSpacing(16)
        root.addLayout(body, 1)

        setup_card = self._card()
        setup_card.setMinimumWidth(400)
        setup_card.setMaximumWidth(440)
        setup_layout = QVBoxLayout(setup_card)
        setup_layout.setContentsMargins(18, 18, 18, 18)
        setup_layout.setSpacing(16)

        setup_title = self._section_title("文件设置", "按顺序选择生成所需的三个文件。")
        setup_layout.addWidget(setup_title)

        setup_layout.addWidget(self._file_picker("Excel 表格", self.excel_edit, "选择 Excel", self.browse_excel))
        setup_layout.addWidget(self._file_picker("PPT 模板", self.ppt_edit, "选择 PPT", self.browse_ppt))
        setup_layout.addWidget(self._file_picker("输出目录", self.output_edit, "选择目录", self.browse_output))

        sheet_row = QWidget()
        sheet_layout = QVBoxLayout(sheet_row)
        sheet_layout.setContentsMargins(0, 0, 0, 0)
        sheet_layout.setSpacing(6)
        sheet_label = QLabel("工作表")
        sheet_label.setObjectName("fieldLabel")
        self.sheet_combo.setMinimumHeight(38)
        self.sheet_combo.currentTextChanged.connect(self.refresh_preview)
        sheet_layout.addWidget(sheet_label)
        sheet_layout.addWidget(self.sheet_combo)
        setup_layout.addWidget(sheet_row)

        self.generate_btn = QPushButton("生成测试 PPT")
        self.generate_btn.setObjectName("primaryButton")
        self.generate_btn.setMinimumHeight(46)
        self.generate_btn.clicked.connect(self.start_generation)
        setup_layout.addWidget(self.generate_btn)

        secondary_actions = QHBoxLayout()
        secondary_actions.setSpacing(10)
        self.reload_btn = QPushButton("重新读取")
        self.reload_btn.setObjectName("secondaryButton")
        self.reload_btn.clicked.connect(self.refresh_preview)
        self.open_output_btn = QPushButton("打开目录")
        self.open_output_btn.setObjectName("secondaryButton")
        self.open_output_btn.clicked.connect(self.open_output_folder)
        secondary_actions.addWidget(self.reload_btn)
        secondary_actions.addWidget(self.open_output_btn)
        setup_layout.addLayout(secondary_actions)
        setup_layout.addStretch(1)

        right = QVBoxLayout()
        right.setSpacing(16)

        summary_card = self._card()
        summary_layout = QVBoxLayout(summary_card)
        summary_layout.setContentsMargins(18, 18, 18, 18)
        summary_layout.setSpacing(14)
        summary_layout.addWidget(self._section_title("读取摘要", "这里显示当前 Excel 的识别结果。"))

        metric_row = QHBoxLayout()
        metric_row.setSpacing(12)
        metric_row.addWidget(self._metric("字段数", self.fields_value))
        metric_row.addWidget(self._metric("数据行", self.rows_value))
        metric_row.addWidget(self._metric("带图行", self.images_value))
        summary_layout.addLayout(metric_row)

        self.progress_bar.setMinimum(0)
        self.progress_bar.setMaximum(100)
        self.progress_bar.setValue(0)
        self.progress_bar.setTextVisible(False)
        self.progress_bar.setMinimumHeight(8)
        summary_layout.addWidget(self.progress_bar)

        log_card = self._card()
        log_layout = QVBoxLayout(log_card)
        log_layout.setContentsMargins(18, 18, 18, 18)
        log_layout.setSpacing(12)
        log_layout.addWidget(self._section_title("运行日志", "生成进度、缺图提示和输出结果会显示在这里。"))
        self.log_box.setReadOnly(True)
        self.log_box.setMinimumHeight(430)
        self.log_box.setPlaceholderText("还没有运行记录")
        log_layout.addWidget(self.log_box, 1)

        right.addWidget(summary_card, 0)
        right.addWidget(log_card, 1)

        body.addWidget(setup_card, 0)
        body.addLayout(right, 1)

        self.statusBar().hide()

        self.setStyleSheet(
            """
            QWidget#root {
                background: #eef2f6;
            }
            QLabel#title {
                font-size: 24px;
                font-weight: 700;
                color: #111827;
            }
            QLabel#subtitle {
                color: #667085;
                font-size: 12px;
            }
            QLabel#sectionTitle {
                font-size: 15px;
                font-weight: 700;
                color: #111827;
            }
            QLabel#sectionHint {
                color: #788394;
                font-size: 12px;
            }
            QLabel#fieldLabel {
                color: #344054;
                font-weight: 600;
            }
            QLabel#statusPill {
                background: #ffffff;
                border: 1px solid #d8dee8;
                border-radius: 14px;
                padding: 6px 12px;
                color: #475467;
            }
            QFrame#card {
                background: #ffffff;
                border: 1px solid #dbe2ec;
                border-radius: 8px;
            }
            QFrame#metric {
                background: #f7f9fc;
                border: 1px solid #e1e7ef;
                border-radius: 6px;
            }
            QLabel#metricLabel {
                color: #667085;
                font-size: 12px;
            }
            QLabel#metricValue {
                color: #101828;
                font-size: 22px;
                font-weight: 700;
            }
            QLineEdit, QComboBox, QPlainTextEdit {
                background: #ffffff;
                border: 1px solid #cfd8e6;
                border-radius: 6px;
                padding: 7px 9px;
                color: #111827;
                selection-background-color: #2563eb;
            }
            QLineEdit:focus, QComboBox:focus, QPlainTextEdit:focus {
                border: 1px solid #2563eb;
            }
            QPushButton {
                border-radius: 6px;
                padding: 8px 12px;
                font-weight: 600;
            }
            QPushButton#primaryButton {
                background: #2563eb;
                color: white;
                border: 1px solid #1d4ed8;
            }
            QPushButton#primaryButton:hover {
                background: #1d4ed8;
            }
            QPushButton#secondaryButton, QPushButton#fileButton {
                background: #ffffff;
                color: #1f2937;
                border: 1px solid #cfd8e6;
            }
            QPushButton#secondaryButton:hover, QPushButton#fileButton:hover {
                background: #f6f8fb;
            }
            QProgressBar {
                border: none;
                background: #e5eaf1;
                border-radius: 4px;
            }
            QProgressBar::chunk {
                background: #2563eb;
                border-radius: 4px;
            }
            """
        )

    def _card(self) -> QFrame:
        frame = QFrame()
        frame.setObjectName("card")
        return frame

    def _section_title(self, title: str, hint: str) -> QWidget:
        widget = QWidget()
        layout = QVBoxLayout(widget)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(3)
        title_label = QLabel(title)
        title_label.setObjectName("sectionTitle")
        hint_label = QLabel(hint)
        hint_label.setObjectName("sectionHint")
        hint_label.setWordWrap(True)
        layout.addWidget(title_label)
        layout.addWidget(hint_label)
        return widget

    def _file_picker(self, title: str, edit: QLineEdit, button_text: str, handler: Callable[[], None]) -> QWidget:
        widget = QWidget()
        layout = QVBoxLayout(widget)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(6)
        label = QLabel(title)
        label.setObjectName("fieldLabel")
        row = QHBoxLayout()
        row.setSpacing(8)
        edit.setMinimumHeight(38)
        edit.setPlaceholderText(title)
        button = QPushButton(button_text)
        button.setObjectName("fileButton")
        button.setMinimumHeight(38)
        button.clicked.connect(handler)
        row.addWidget(edit, 1)
        row.addWidget(button, 0)
        layout.addWidget(label)
        layout.addLayout(row)
        return widget

    def _metric(self, label: str, value_label: QLabel) -> QFrame:
        frame = QFrame()
        frame.setObjectName("metric")
        layout = QVBoxLayout(frame)
        layout.setContentsMargins(14, 12, 14, 12)
        layout.setSpacing(3)
        name = QLabel(label)
        name.setObjectName("metricLabel")
        value_label.setObjectName("metricValue")
        layout.addWidget(name)
        layout.addWidget(value_label)
        return frame

    def _load_defaults(self) -> None:
        if DEFAULT_EXCEL:
            self.excel_edit.setText(display_path(DEFAULT_EXCEL))
            self.excel_edit.setToolTip(str(DEFAULT_EXCEL.resolve()))
        if DEFAULT_PPT:
            self.ppt_edit.setText(display_path(DEFAULT_PPT))
            self.ppt_edit.setToolTip(str(DEFAULT_PPT.resolve()))
        self.output_edit.setText(display_path(DEFAULT_OUTPUT))
        self.output_edit.setToolTip(str(DEFAULT_OUTPUT.resolve()))
        if self.excel_edit.text():
            self.refresh_preview()

    def append_log(self, text: str) -> None:
        self.log_box.appendPlainText(text)

    def browse_excel(self) -> None:
        path, _ = QFileDialog.getOpenFileName(self, "选择 Excel 文件", str(ROOT), "Excel Files (*.xlsx *.xlsm)")
        if path:
            p = Path(path)
            self.excel_edit.setText(display_path(p))
            self.excel_edit.setToolTip(str(p.resolve()))
            self.refresh_preview()

    def browse_ppt(self) -> None:
        path, _ = QFileDialog.getOpenFileName(self, "选择 PPT 模板", str(ROOT), "PowerPoint Files (*.pptx)")
        if path:
            p = Path(path)
            self.ppt_edit.setText(display_path(p))
            self.ppt_edit.setToolTip(str(p.resolve()))

    def browse_output(self) -> None:
        path = QFileDialog.getExistingDirectory(
            self,
            "选择输出目录",
            str(resolve_path_text(self.output_edit.text()) if self.output_edit.text() else DEFAULT_OUTPUT_DIR),
        )
        if path:
            p = Path(path)
            self.output_edit.setText(display_path(p))
            self.output_edit.setToolTip(str(p.resolve()))

    def refresh_preview(self) -> None:
        excel_path = resolve_path_text(self.excel_edit.text())
        if not excel_path.exists():
            self.ready_label.setText("等待 Excel")
            self.fields_value.setText("-")
            self.rows_value.setText("-")
            self.images_value.setText("-")
            self.sheet_combo.blockSignals(True)
            self.sheet_combo.clear()
            self.sheet_combo.blockSignals(False)
            return

        try:
            wb = load_workbook(excel_path, data_only=True)
            sheets = wb.sheetnames
            self.sheet_combo.blockSignals(True)
            current = self.sheet_combo.currentText()
            self.sheet_combo.clear()
            self.sheet_combo.addItems(sheets)
            if current in sheets:
                self.sheet_combo.setCurrentText(current)
            self.sheet_combo.blockSignals(False)
            target_sheet = self.sheet_combo.currentText() or sheets[0]
            headers, records, stats = extract_excel_rows_and_images(excel_path, target_sheet, TEMP_DIR)
            self.headers = headers
            self.records = records
            self.stats = stats
            self._update_summary()
            self.ready_label.setText("已读取")
            self.append_log(f"已读取: {excel_path.name} / {target_sheet}")
        except Exception as exc:
            self.ready_label.setText("读取失败")
            self.append_log(traceback.format_exc())
            QMessageBox.warning(self, "读取失败", str(exc))

    def _update_summary(self) -> None:
        if not self.records:
            self.fields_value.setText("-")
            self.rows_value.setText("0")
            self.images_value.setText("0")
            return
        self.fields_value.setText(str(len(self.headers)))
        self.rows_value.setText(str(self.stats.get("rows", 0)))
        self.images_value.setText(str(self.stats.get("rows_with_image", 0)))

    def start_generation(self) -> None:
        excel_path = resolve_path_text(self.excel_edit.text())
        ppt_path = resolve_path_text(self.ppt_edit.text())
        output_path = resolve_output_target(self.output_edit.text())
        sheet_name = self.sheet_combo.currentText().strip()

        if not excel_path.exists():
            QMessageBox.warning(self, "提示", "请先选择有效的 Excel 文件")
            return
        if not ppt_path.exists():
            QMessageBox.warning(self, "提示", "请先选择有效的 PPT 模板")
            return
        if not sheet_name:
            QMessageBox.warning(self, "提示", "请先选择工作表")
            return

        self.generate_btn.setEnabled(False)
        self.progress_bar.setValue(0)
        self.ready_label.setText("生成中")
        self.append_log("开始生成...")
        self.worker = GenerateWorker(excel_path, ppt_path, output_path, sheet_name)
        self.worker.log_signal.connect(self.append_log)
        self.worker.progress_signal.connect(self.on_progress)
        self.worker.done_signal.connect(self.on_done)
        self.worker.error_signal.connect(self.on_error)
        self.worker.finished.connect(lambda: self.generate_btn.setEnabled(True))
        self.worker.start()

    def on_progress(self, cur: int, total: int) -> None:
        value = int(cur * 100 / total) if total else 0
        self.progress_bar.setValue(value)

    def on_done(self, output_file: str) -> None:
        self.progress_bar.setValue(100)
        self.ready_label.setText("生成完成")
        self.append_log(f"完成: {output_file}")
        QMessageBox.information(self, "完成", f"PPT 已生成:\n{output_file}")

    def on_error(self, error_text: str) -> None:
        self.progress_bar.setValue(0)
        self.ready_label.setText("生成失败")
        self.append_log(error_text)
        if "PermissionError" in error_text or "Permission denied" in error_text:
            message = (
                "输出文件正在被占用，无法覆盖保存。\n\n"
                "请关闭已经打开的 PPT 文件，或者换一个输出文件名后再生成。"
            )
        else:
            message = "生成失败，详细错误已经写入日志。"
        QMessageBox.critical(self, "生成失败", message)

    def open_output_folder(self) -> None:
        output_path = resolve_output_target(self.output_edit.text() or str(DEFAULT_OUTPUT_DIR))
        folder = output_path.parent
        folder.mkdir(parents=True, exist_ok=True)
        try:
            os.startfile(str(folder))
        except Exception as exc:
            QMessageBox.warning(self, "提示", f"无法打开输出目录: {exc}")


def run_headless(args: argparse.Namespace) -> int:
    excel_path = Path(args.excel or (DEFAULT_EXCEL or ""))
    ppt_path = Path(args.ppt or (DEFAULT_PPT or ""))
    output_path = resolve_output_target(args.output or str(DEFAULT_OUTPUT_DIR))
    if not excel_path.exists():
        print(f"Excel 文件不存在: {excel_path}")
        return 1
    if not ppt_path.exists():
        print(f"PPT 模板不存在: {ppt_path}")
        return 1
    try:
        _, saved_path = generate_ppt(
            excel_path,
            ppt_path,
            output_path,
            args.sheet,
            log=print,
            progress=lambda cur, total: print(f"进度: {cur}/{total}"),
        )
        print(f"完成: {saved_path}")
        return 0
    except Exception:
        print(traceback.format_exc())
        return 1


def main() -> int:
    parser = argparse.ArgumentParser(description="Excel 批量生成 PPT 测试版")
    parser.add_argument("--headless", action="store_true", help="不启动界面，直接生成")
    parser.add_argument("--excel", type=str, help="Excel 文件路径")
    parser.add_argument("--ppt", type=str, help="PPT 模板路径")
    parser.add_argument("--output", type=str, help="输出 PPT 路径")
    parser.add_argument("--sheet", type=str, help="工作表名称")
    args = parser.parse_args()

    if args.headless:
        return run_headless(args)

    app = QApplication(sys.argv)
    win = CleanMainWindow()
    win.show()
    return app.exec_()


if __name__ == "__main__":
    raise SystemExit(main())
